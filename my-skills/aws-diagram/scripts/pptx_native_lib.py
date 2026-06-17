"""
Native PPTX builder library for AWS architecture diagrams.

Battle-tested helper layer for building fully-editable PowerPoint
architecture slides with python-pptx. Born from a field session,
where the automated px->EMU exporter (pptx_export.py) produced
20 visual defects (container/icon overlap, orphaned labels, arrows
piercing icons) and a hand-written script with explicit inch coordinates
produced a clean slide on the first full pass.

Design principles (the "why" behind each helper):

1. Explicit inch coordinates, no px->EMU translation.
   The SVG layout engine reasons about icon boxes but not about rendered
   text extents (especially CJK text, which is ~1.7x wider per char than
   the Latin metrics python-pptx assumes). Translating those coords to
   EMU inherits every blind spot. Placing shapes directly in inches keeps
   the coordinate system human-debuggable: a defect at "x=8.65" is
   immediately actionable, a defect at "EMU 7920757" is not.

2. word_wrap=False + generous fixed-width label boxes, centered on icon.
   PowerPoint wraps CJK labels aggressively in auto-sized boxes
   ("AgentCore Runtime" -> 3 broken lines). Fixed 1.4" boxes centered on
   the icon midpoint render one clean line in PowerPoint, Keynote, and
   LibreOffice alike.

3. Strict z-order: containers -> lines -> labels -> icons+labels.
   python-pptx stacks shapes in insertion order. Drawing containers
   first and nodes last means an arrow can pass "under" an icon
   gracefully instead of slicing through it.

4. Arrow labels get a white fill.
   A label sitting on a dashed line is unreadable; the white box masks
   the line segment behind the text, which reads as a deliberate label
   gap (the same trick the AWS icon deck uses).

5. Multi-segment routes are explicit, not solved.
   Orthogonal auto-routing is exactly what failed in the exporter. Here
   an L- or Z-shaped route is just 2-3 straight line() calls whose
   coordinates you choose. The render-QA loop (see
   references/pptx-native-workflow.md) catches the mistakes cheaply.

Usage:
    from pptx_native_lib import NativeSlideBuilder, DARK, PURPLE

    b = NativeSlideBuilder()                      # new 16:9 deck
    b.title("My Architecture", "subtitle text")
    b.container(0.5, 1.5, 4.0, 5.0, "AWS Cloud", dark_header=True)
    b.line(1.5, 3.0, 3.0, 3.0, tail=True)         # arrow ->
    b.arrow_label(2.25, 2.8, "HTTPS")
    b.node(1.5, 3.0, "lambda", "AWS Lambda", "handler")
    b.save("out.pptx")

Icons: pass icon_dir pointing at rasterized PNGs (icon name -> PNG file).
Rasterize once per project:
    rsvg-convert -w 220 -h 220 icons/lambda.svg -o icon-png/lambda.png
"""

from __future__ import annotations

import subprocess
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.oxml.ns import qn

# ── AWS design-system colors ──────────────────────────────────────────
DARK = RGBColor(0x23, 0x2F, 0x3E)      # Squid Ink — outer borders, headers
GRAY = RGBColor(0x5A, 0x6B, 0x7B)      # sublabels, container labels
BORDER = RGBColor(0xAE, 0xB6, 0xBF)    # generic container border
ARROW = RGBColor(0x54, 0x5B, 0x64)     # default arrow gray
PURPLE = RGBColor(0x8C, 0x4F, 0xFF)    # networking / special flows
TEAL = RGBColor(0x01, 0xA8, 0x8D)      # AI/ML accent
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

ICON_IN = 0.62  # icon edge length in inches; labels assume this size


class NativeSlideBuilder:
    """Builds one native-shape architecture slide on a 13.333x7.5 canvas.

    All coordinates are inches. (cx, cy) always means the CENTER of an
    icon; container/line calls take absolute edge coordinates.
    """

    def __init__(self, prs=None, icon_dir: str | Path = "icon-png"):
        if prs is None:
            prs = Presentation()
            prs.slide_width = Inches(13.333)
            prs.slide_height = Inches(7.5)
        self.prs = prs
        self.slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
        self.shapes = self.slide.shapes
        self.icon_dir = Path(icon_dir)

    # ── canvas chrome ─────────────────────────────────────────────────

    def title(self, text: str, subtitle: str | None = None) -> None:
        t = self.shapes.add_textbox(Inches(0.45), Inches(0.18),
                                    Inches(12.0), Inches(0.4))
        tf = t.text_frame
        tf.word_wrap = False
        r = tf.paragraphs[0].add_run()
        r.text = text
        r.font.size = Pt(18)
        r.font.bold = True
        r.font.color.rgb = DARK
        if subtitle:
            t2 = self.shapes.add_textbox(Inches(0.47), Inches(0.58),
                                         Inches(12.0), Inches(0.25))
            tf2 = t2.text_frame
            tf2.word_wrap = False
            r2 = tf2.paragraphs[0].add_run()
            r2.text = subtitle
            r2.font.size = Pt(10)
            r2.font.color.rgb = GRAY

    # ── containers (draw these FIRST — z-order bottom) ────────────────

    def container(self, x1: float, y1: float, x2: float, y2: float,
                  label: str, border_color: RGBColor = BORDER,
                  width_pt: float = 1.0, dark_header: bool = False):
        """Rounded-rect grouping box from (x1,y1) to (x2,y2).

        dark_header=True draws the AWS-Cloud-style dark title bar; use it
        only for the outermost account/cloud boundary. Plain containers
        get a small gray label inside the top-left corner.
        Keep >= 0.25in clearance between a container edge and any icon
        center +- ICON_IN/2 that does not belong to it.
        """
        box = self.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                    Inches(x1), Inches(y1),
                                    Inches(x2 - x1), Inches(y2 - y1))
        box.adjustments[0] = 0.03
        box.fill.background()
        box.line.color.rgb = border_color
        box.line.width = Pt(width_pt)
        box.shadow.inherit = False
        if dark_header:
            hdr = self.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x1),
                                        Inches(y1), Inches(1.45), Inches(0.26))
            hdr.fill.solid()
            hdr.fill.fore_color.rgb = DARK
            hdr.line.fill.background()
            hdr.shadow.inherit = False
            tf = hdr.text_frame
            tf.margin_left = Inches(0.06)
            tf.margin_top = Inches(0.01)
            tf.word_wrap = False
            r = tf.paragraphs[0].add_run()
            r.text = label
            r.font.size = Pt(10)
            r.font.bold = True
            r.font.color.rgb = WHITE
        else:
            t = self.shapes.add_textbox(Inches(x1 + 0.06), Inches(y1 + 0.02),
                                        Inches(2.2), Inches(0.2))
            tf = t.text_frame
            tf.margin_left = Emu(0)
            tf.margin_top = Emu(0)
            tf.word_wrap = False
            r = tf.paragraphs[0].add_run()
            r.text = label
            r.font.size = Pt(9)
            r.font.color.rgb = GRAY
        return box

    # ── connectors (draw AFTER containers, BEFORE nodes) ──────────────

    def line(self, x1: float, y1: float, x2: float, y2: float,
             dashed: bool = False, color: RGBColor = ARROW,
             head: bool = False, tail: bool = False,
             width_pt: float = 1.5):
        """One straight connector segment. tail=True puts an arrowhead at
        (x2,y2); head=True at (x1,y1); both for bidirectional.

        Build L/Z routes from 2-3 calls and put the arrowhead only on the
        final segment. Route vertical drops through empty corridors —
        when a route must pass an icon column, jog around it rather than
        through it. Arrowheads should land ~0.05-0.10in short of the icon
        edge (icon edge = center +- ICON_IN/2), and a segment entering
        from below must stop BELOW the label block (center + ICON_IN/2 +
        ~0.5in), or it will spear the text.
        """
        conn = self.shapes.add_connector(MSO_CONNECTOR.STRAIGHT,
                                         Inches(x1), Inches(y1),
                                         Inches(x2), Inches(y2))
        conn.line.color.rgb = color
        conn.line.width = Pt(width_pt)
        conn.shadow.inherit = False
        ln = conn.line._get_or_add_ln()
        if dashed:
            ln.append(ln.makeelement(qn('a:prstDash'), {'val': 'dash'}))
        if tail:
            ln.append(ln.makeelement(qn('a:tailEnd'),
                                     {'type': 'triangle', 'w': 'med', 'len': 'med'}))
        if head:
            ln.append(ln.makeelement(qn('a:headEnd'),
                                     {'type': 'triangle', 'w': 'med', 'len': 'med'}))
        return conn

    def arrow_label(self, cx: float, cy: float, text: str,
                    color: RGBColor = DARK, size: int = 9,
                    w: float = 1.6):
        """Italic flow label centered at (cx,cy), white-filled so it masks
        the line behind it. Place ~0.22in above the line midpoint, and
        size w to the text (~0.11in per Korean char, ~0.07 per Latin) so
        the white box doesn't blank out more line than needed."""
        t = self.shapes.add_textbox(Inches(cx - w / 2), Inches(cy - 0.11),
                                    Inches(w), Inches(0.22))
        t.fill.solid()
        t.fill.fore_color.rgb = WHITE
        tf = t.text_frame
        tf.margin_left = Emu(0)
        tf.margin_top = Emu(0)
        tf.margin_right = Emu(0)
        tf.margin_bottom = Emu(0)
        tf.word_wrap = False
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = text
        r.font.size = Pt(size)
        r.font.italic = True
        r.font.color.rgb = color
        return t

    # ── nodes (draw LAST — z-order top) ───────────────────────────────

    def node(self, cx: float, cy: float, icon: str, name: str,
             sub: str | None = None):
        """Service icon + 2-line label centered at (cx,cy).

        Label box is fixed 1.4in wide with word_wrap=False — this is the
        fix for CJK label wrapping. Keep horizontal neighbors >= 1.55in
        apart (center-to-center) so 1.4in label boxes never collide.
        Vertical pitch >= 1.25in keeps a label clear of the icon below.
        """
        png = self.icon_dir / f"{icon}.png"
        if not png.exists():
            raise FileNotFoundError(
                f"{png} not found — rasterize it first, e.g.\n"
                f"  rsvg-convert -w 220 -h 220 icons/{icon}.svg -o {png}")
        half = ICON_IN / 2
        self.shapes.add_picture(str(png), Inches(cx - half), Inches(cy - half),
                                Inches(ICON_IN), Inches(ICON_IN))
        t = self.shapes.add_textbox(Inches(cx - 0.7), Inches(cy + half + 0.03),
                                    Inches(1.4), Inches(0.5))
        tf = t.text_frame
        tf.margin_left = Emu(0)
        tf.margin_top = Emu(0)
        tf.margin_right = Emu(0)
        tf.margin_bottom = Emu(0)
        tf.word_wrap = False
        p1 = tf.paragraphs[0]
        p1.alignment = PP_ALIGN.CENTER
        r1 = p1.add_run()
        r1.text = name
        r1.font.size = Pt(10)
        r1.font.color.rgb = DARK
        if sub:
            p2 = tf.add_paragraph()
            p2.alignment = PP_ALIGN.CENTER
            r2 = p2.add_run()
            r2.text = sub
            r2.font.size = Pt(8)
            r2.font.italic = True
            r2.font.color.rgb = GRAY

    # ── output ────────────────────────────────────────────────────────

    def save(self, path: str | Path) -> None:
        self.prs.save(str(path))


def rasterize_icons(icon_names: list[str], svg_dir: str | Path,
                    out_dir: str | Path = "icon-png", px: int = 220) -> None:
    """Rasterize the named skill icons to PNG via rsvg-convert."""
    svg_dir, out_dir = Path(svg_dir), Path(out_dir)
    out_dir.mkdir(exist_ok=True)
    for name in icon_names:
        src = svg_dir / f"{name}.svg"
        if not src.exists():
            raise FileNotFoundError(f"icon not found: {src}")
        subprocess.run(["rsvg-convert", "-w", str(px), "-h", str(px),
                        str(src), "-o", str(out_dir / f"{name}.png")],
                       check=True)

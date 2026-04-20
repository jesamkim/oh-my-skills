"""
Native PPTX export for AWS architecture diagrams.

Renders each diagram element as an individual, editable PowerPoint shape:
- Service icons: Picture shapes (SVG rasterized to PNG)
- Arrows: Native connector shapes via OOXML injection
- Containers: Rectangle/RoundedRectangle shapes
- Labels: TextBox shapes (editable text)
- Step callouts: Oval shapes with centered numbers

Produces LAYOUT_WIDE (13.33" x 7.50") slides matching the AWS Architecture
Icon Deck design system.
"""

from __future__ import annotations

import io
from pathlib import Path
from typing import Optional

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor

from diagram_schema import DiagramDefinition
from layout_engine import LayoutResult, ICON_SIZE
from orthogonal_router import RoutedPath, CALLOUT_RADIUS
from svg_renderer import THEME_COLORS, CONTAINER_STYLES


# ============================================================
# Constants
# ============================================================

EMU_PER_INCH = 914400
SLIDE_WIDTH_IN = 13.33
SLIDE_HEIGHT_IN = 7.5

MARGIN_TOP_IN = 0.8
MARGIN_SIDES_IN = 0.3
MARGIN_BOTTOM_IN = 0.5

USABLE_WIDTH_EMU = int((SLIDE_WIDTH_IN - 2 * MARGIN_SIDES_IN) * EMU_PER_INCH)
USABLE_HEIGHT_EMU = int((SLIDE_HEIGHT_IN - MARGIN_TOP_IN - MARGIN_BOTTOM_IN) * EMU_PER_INCH)

FONT_NAME = "Calibri"

# AWS Arrow spec
ARROW_WIDTH_EMU = 12700   # 1.0pt
ARROW_WIDTH_EMPHASIS_EMU = 15875  # 1.2pt


# ============================================================
# Coordinate Translation
# ============================================================

def _compute_scale_and_offset(
    viewbox_w: int, viewbox_h: int,
) -> tuple[float, int, int]:
    """Compute scale factor and centering offset for pixel->EMU conversion.

    Returns (scale, offset_x_emu, offset_y_emu).
    """
    scale_x = USABLE_WIDTH_EMU / max(viewbox_w, 1)
    scale_y = USABLE_HEIGHT_EMU / max(viewbox_h, 1)
    scale = min(scale_x, scale_y)

    content_w = int(viewbox_w * scale)
    content_h = int(viewbox_h * scale)

    offset_x = int(MARGIN_SIDES_IN * EMU_PER_INCH) + (USABLE_WIDTH_EMU - content_w) // 2
    offset_y = int(MARGIN_TOP_IN * EMU_PER_INCH) + (USABLE_HEIGHT_EMU - content_h) // 2

    return scale, offset_x, offset_y


def _px_to_emu(px: float, scale: float) -> int:
    """Convert pixel value to EMU."""
    return int(px * scale)


def _px_xy(px_x: float, px_y: float, scale: float, ox: int, oy: int) -> tuple[int, int]:
    """Convert pixel coordinates to slide EMU coordinates."""
    return (int(px_x * scale) + ox, int(px_y * scale) + oy)


# ============================================================
# Color Helper
# ============================================================

def _rgb(hex_color: str) -> RGBColor:
    """Convert hex color string (without #) to RGBColor."""
    h = hex_color.lstrip("#")
    return RGBColor.from_string(h)


# ============================================================
# Container Rendering
# ============================================================

def _add_containers(slide, layout: LayoutResult, diagram: DiagramDefinition,
                    scale: float, ox: int, oy: int):
    """Add container rectangles to the slide."""
    theme_name = diagram.theme if diagram.theme in CONTAINER_STYLES else "light"
    styles = CONTAINER_STYLES[theme_name]

    for cl in layout.containers:
        style = styles.get(cl.container_type, styles["generic"])
        b = cl.bbox

        left, top = _px_xy(b.x, b.y, scale, ox, oy)
        width = _px_to_emu(b.w, scale)
        height = _px_to_emu(b.h, scale)

        shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if style.get("rx", "0") not in ("0", None, "") else MSO_SHAPE.RECTANGLE
        shape = slide.shapes.add_shape(shape_type, left, top, width, height)

        # Fill
        fill_color = style.get("fill", "none")
        if fill_color == "none" or not fill_color:
            shape.fill.background()
        else:
            shape.fill.solid()
            shape.fill.fore_color.rgb = _rgb(fill_color)
            opacity = style.get("fill-opacity")
            if opacity:
                from lxml import etree
                alpha_val = str(int(float(opacity) * 100000))
                ns = "http://schemas.openxmlformats.org/drawingml/2006/main"
                # Access the underlying XML element for the fill
                sp_elem = shape._element
                for clr in sp_elem.iter(f"{{{ns}}}srgbClr"):
                    existing = clr.find(f"{{{ns}}}alpha")
                    if existing is None:
                        alpha_elem = etree.SubElement(clr, f"{{{ns}}}alpha")
                        alpha_elem.set("val", alpha_val)

        # Border
        shape.line.color.rgb = _rgb(style["stroke"])
        shape.line.width = Pt(float(style["stroke-width"]))

        # Dash style
        dash = style.get("stroke-dasharray")
        if dash:
            from pptx.enum.dml import MSO_LINE_DASH_STYLE
            shape.line.dash_style = MSO_LINE_DASH_STYLE.DASH

        # No text on the shape itself
        if shape.has_text_frame:
            shape.text_frame.clear()

        # Header
        if style.get("header_bg"):
            header_w = min(int(len(cl.label) * 8 + 40), int(b.w))
            h_left, h_top = left, top
            h_width = _px_to_emu(header_w, scale)
            h_height = _px_to_emu(22, scale)

            header_shape = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE, h_left, h_top, h_width, h_height
            )
            header_shape.fill.solid()
            header_shape.fill.fore_color.rgb = _rgb(style["header_bg"])
            header_shape.line.fill.background()

            tf = header_shape.text_frame
            tf.word_wrap = False
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.LEFT
            run = p.add_run()
            run.text = cl.label
            run.font.size = Pt(9)
            run.font.bold = True
            run.font.name = FONT_NAME
            run.font.color.rgb = _rgb(style["header_text"])
            # Vertical centering
            tf.auto_size = None
            tf.margin_left = Emu(_px_to_emu(6, scale))
            tf.margin_top = Emu(0)
        else:
            # Simple text label
            lbl_left = left + _px_to_emu(8, scale)
            lbl_top = top + _px_to_emu(2, scale)
            lbl_w = _px_to_emu(len(cl.label) * 8, scale)
            lbl_h = _px_to_emu(16, scale)

            txbox = slide.shapes.add_textbox(lbl_left, lbl_top, lbl_w, lbl_h)
            p = txbox.text_frame.paragraphs[0]
            run = p.add_run()
            run.text = cl.label
            run.font.size = Pt(8)
            run.font.name = FONT_NAME
            run.font.color.rgb = _rgb(style["header_text"])


# ============================================================
# Icon Placement
# ============================================================

def _add_icons(slide, layout: LayoutResult, diagram: DiagramDefinition,
               icons_dir: Path, scale: float, ox: int, oy: int):
    """Add service icons as individual Picture shapes."""
    from icon_rasterizer import rasterize_icon, find_icon_path

    for node in diagram.nodes:
        nl = layout.nodes.get(node.id)
        if not nl:
            continue

        icon_path = find_icon_path(node.service, icons_dir)
        left, top = _px_xy(nl.icon_x, nl.icon_y, scale, ox, oy)
        icon_emu = _px_to_emu(ICON_SIZE, scale)

        if icon_path:
            try:
                png_bytes = rasterize_icon(icon_path, size=128)
                pic = slide.shapes.add_picture(
                    io.BytesIO(png_bytes), left, top, icon_emu, icon_emu
                )
                continue
            except Exception:
                pass

        # Fallback: colored rectangle with service name
        shape = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, left, top, icon_emu, icon_emu
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = _rgb("545B64")
        shape.line.fill.background()
        tf = shape.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = node.service
        run.font.size = Pt(7)
        run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        run.font.name = FONT_NAME


# ============================================================
# Label Rendering
# ============================================================

def _add_node_labels(slide, layout: LayoutResult, diagram: DiagramDefinition,
                     scale: float, ox: int, oy: int):
    """Add node labels as editable TextBox shapes."""
    theme_name = diagram.theme if diagram.theme in THEME_COLORS else "light"
    theme = THEME_COLORS[theme_name]

    for node in diagram.nodes:
        nl = layout.nodes.get(node.id)
        if not nl:
            continue

        label_lines = node.label.split("\n") if node.label else []
        num_lines = max(len(label_lines), 1)
        sublabel_lines = 1 if node.sublabel else 0

        box_w = _px_to_emu(100, scale)
        box_h = _px_to_emu((num_lines + sublabel_lines) * 14 + 4, scale)
        center_x, _ = _px_xy(nl.label_x, nl.label_y, scale, ox, oy)
        _, label_top = _px_xy(0, nl.label_y, scale, 0, oy)
        label_left = center_x - box_w // 2

        txbox = slide.shapes.add_textbox(label_left, label_top, box_w, box_h)
        tf = txbox.text_frame
        tf.word_wrap = True
        tf.auto_size = None

        # Primary label
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        text = "\n".join(label_lines) if label_lines else ""
        run = p.add_run()
        run.text = text
        run.font.size = Pt(9)
        run.font.name = FONT_NAME
        run.font.color.rgb = _rgb(theme["text"])

        # Sublabel
        if node.sublabel:
            p2 = tf.add_paragraph()
            p2.alignment = PP_ALIGN.CENTER
            run2 = p2.add_run()
            run2.text = node.sublabel
            run2.font.size = Pt(8)
            run2.font.italic = True
            run2.font.name = FONT_NAME
            run2.font.color.rgb = _rgb(theme["label"])


def _add_connection_labels(slide, routes: list[RoutedPath],
                           theme: dict, scale: float, ox: int, oy: int):
    """Add arrow labels as TextBox shapes."""
    for route in routes:
        if not route.label or not route.label_pos:
            continue

        lx, ly = route.label_pos
        center_x, label_y = _px_xy(lx, ly, scale, ox, oy)
        box_w = _px_to_emu(80, scale)
        box_h = _px_to_emu(16, scale)

        txbox = slide.shapes.add_textbox(
            center_x - box_w // 2, label_y - box_h // 2, box_w, box_h
        )
        tf = txbox.text_frame
        tf.word_wrap = False
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = route.label
        run.font.size = Pt(8)
        run.font.italic = True
        run.font.name = FONT_NAME
        run.font.color.rgb = _rgb(theme["text"])


# ============================================================
# Arrows (Native Connectors)
# ============================================================

def _add_arrows(slide, routes: list[RoutedPath],
                theme_name: str, scale: float, ox: int, oy: int):
    """Add arrows as native PPTX shapes."""
    from pptx_connector import add_straight_connector, add_freeform_arrow

    theme = THEME_COLORS.get(theme_name, THEME_COLORS["light"])
    arrow_color = theme["arrow"].lstrip("#")

    for route in routes:
        if len(route.waypoints) < 2:
            continue

        # Convert waypoints to EMU
        waypoints_emu = [
            _px_xy(wx, wy, scale, ox, oy)
            for wx, wy in route.waypoints
        ]

        style_config = {
            "color": arrow_color,
            "width_emu": ARROW_WIDTH_EMU,
            "dash": route.style == "dashed",
            "bidirectional": route.style == "bidirectional",
        }

        if len(waypoints_emu) == 2:
            add_straight_connector(
                slide,
                waypoints_emu[0][0], waypoints_emu[0][1],
                waypoints_emu[1][0], waypoints_emu[1][1],
                style_config,
            )
        else:
            add_freeform_arrow(slide, waypoints_emu, style_config)


# ============================================================
# Step Number Callouts
# ============================================================

def _add_callouts(slide, routes: list[RoutedPath],
                  theme_name: str, scale: float, ox: int, oy: int):
    """Add step number circles as Oval shapes."""
    theme = THEME_COLORS.get(theme_name, THEME_COLORS["light"])

    if theme_name == "dark":
        fill_color = "FFFFFF"
        text_color = "232F3E"
    else:
        fill_color = "232F3E"
        text_color = "FFFFFF"

    diameter_px = CALLOUT_RADIUS * 2
    diameter_emu = _px_to_emu(diameter_px, scale)

    for route in routes:
        if route.step_number is None or not route.step_pos:
            continue

        sx, sy = route.step_pos
        center_x, center_y = _px_xy(sx, sy, scale, ox, oy)
        left = center_x - diameter_emu // 2
        top = center_y - diameter_emu // 2

        shape = slide.shapes.add_shape(
            MSO_SHAPE.OVAL, left, top, diameter_emu, diameter_emu
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = _rgb(fill_color)
        shape.line.fill.background()

        tf = shape.text_frame
        tf.auto_size = None
        tf.margin_left = Emu(0)
        tf.margin_right = Emu(0)
        tf.margin_top = Emu(0)
        tf.margin_bottom = Emu(0)

        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = str(route.step_number)
        run.font.size = Pt(9)
        run.font.bold = True
        run.font.name = FONT_NAME
        run.font.color.rgb = _rgb(text_color)

        # Vertical center
        tf.auto_size = None
        from pptx.enum.text import MSO_ANCHOR
        tf.paragraphs[0].space_before = Pt(0)
        tf.paragraphs[0].space_after = Pt(0)


# ============================================================
# Title Slide
# ============================================================

def _add_title_slide(prs: Presentation, title: str, subtitle: Optional[str]):
    """Add a white title slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    txbox = slide.shapes.add_textbox(
        Inches(0.5), Inches(2.5), Inches(12.0), Inches(2.0)
    )
    tf = txbox.text_frame
    tf.word_wrap = True

    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = title
    run.font.size = Pt(36)
    run.font.bold = True
    run.font.name = FONT_NAME
    p.alignment = PP_ALIGN.LEFT

    if subtitle:
        p2 = tf.add_paragraph()
        run2 = p2.add_run()
        run2.text = subtitle
        run2.font.size = Pt(18)
        run2.font.name = FONT_NAME
        run2.font.color.rgb = _rgb("545B64")
        p2.alignment = PP_ALIGN.LEFT
        p2.space_before = Pt(12)


# ============================================================
# Main Export
# ============================================================

def export_pptx_native(
    diagram: DiagramDefinition,
    layout: LayoutResult,
    routes: list[RoutedPath],
    icons_dir: str,
    output_path: str,
) -> None:
    """Export a diagram as a native-shape PPTX presentation.

    Every element is an individual, editable PowerPoint shape:
    icons as pictures, arrows as connectors, containers as rectangles,
    labels as textboxes, callouts as ovals.
    """
    prs = Presentation()
    prs.slide_width = Inches(SLIDE_WIDTH_IN)
    prs.slide_height = Inches(SLIDE_HEIGHT_IN)

    # Title slide
    _add_title_slide(prs, diagram.title, diagram.subtitle)

    # Diagram slide
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Title text on diagram slide
    txbox = slide.shapes.add_textbox(
        Inches(0.3), Inches(0.1), Inches(12.0), Inches(0.5)
    )
    p = txbox.text_frame.paragraphs[0]
    run = p.add_run()
    run.text = diagram.title
    run.font.size = Pt(18)
    run.font.bold = True
    run.font.name = FONT_NAME
    if diagram.subtitle:
        run2 = p.add_run()
        run2.text = f"  {diagram.subtitle}"
        run2.font.size = Pt(10)
        run2.font.name = FONT_NAME
        run2.font.color.rgb = _rgb("666666")

    # Compute scale
    scale, ox, oy = _compute_scale_and_offset(
        layout.viewbox_width, layout.viewbox_height
    )

    theme_name = diagram.theme if diagram.theme in THEME_COLORS else "light"
    theme = THEME_COLORS[theme_name]
    icons_path = Path(icons_dir)

    # Set slide background color based on theme
    bg = slide.background
    bg.fill.solid()
    bg.fill.fore_color.rgb = _rgb(theme["background"])

    # Render layers (back to front)
    _add_containers(slide, layout, diagram, scale, ox, oy)
    _add_icons(slide, layout, diagram, icons_path, scale, ox, oy)
    _add_node_labels(slide, layout, diagram, scale, ox, oy)
    _add_arrows(slide, routes, theme_name, scale, ox, oy)
    _add_connection_labels(slide, routes, theme, scale, ox, oy)
    _add_callouts(slide, routes, theme_name, scale, ox, oy)

    # Save
    Path(output_path).parent.mkdir(parents=True, exist_ok=True)
    prs.save(output_path)


# ============================================================
# Legacy API (backward compatibility)
# ============================================================

def export_pptx(
    png_path: str,
    output_path: str,
    title: str = "AWS Architecture Diagram",
    subtitle: Optional[str] = None,
) -> None:
    """Legacy export: embed a PNG diagram into a PPTX slide.

    Deprecated: Use export_pptx_native() for editable native shapes.
    """
    from PIL import Image as PILImage

    prs = Presentation()
    prs.slide_width = Inches(SLIDE_WIDTH_IN)
    prs.slide_height = Inches(SLIDE_HEIGHT_IN)

    _add_title_slide(prs, title, subtitle)

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    txbox = slide.shapes.add_textbox(
        Inches(0.3), Inches(0.15), Inches(12.0), Inches(0.5)
    )
    p = txbox.text_frame.paragraphs[0]
    run = p.add_run()
    run.text = title
    run.font.size = Pt(20)
    run.font.bold = True
    run.font.name = FONT_NAME

    img = PILImage.open(png_path)
    img_w, img_h = img.size
    img.close()

    max_w = Inches(SLIDE_WIDTH_IN - 2 * MARGIN_SIDES_IN)
    max_h = Inches(SLIDE_HEIGHT_IN - MARGIN_TOP_IN - MARGIN_BOTTOM_IN)
    aspect = img_w / img_h

    if aspect > (max_w / max_h):
        width = max_w
        height = int(max_w / aspect)
    else:
        height = max_h
        width = int(max_h * aspect)

    left = (Inches(SLIDE_WIDTH_IN) - width) // 2
    top = Inches(MARGIN_TOP_IN)
    slide.shapes.add_picture(png_path, left, top, width, height)

    Path(output_path).parent.mkdir(parents=True, exist_ok=True)
    prs.save(output_path)

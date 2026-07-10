#!/usr/bin/env python3
"""Extract text from PPTX at run-level for format-preserving translation.

Usage:
    python extract_text.py input.pptx [--per-slide] [--output-dir DIR] [--target-lang LANG]

Output:
    Default:     {filename}_text.json (single file)
    --per-slide: {filename}_slide_{N}.json (one per slide)
"""

import argparse
import json
import os
import sys
from pathlib import Path

try:
    from pptx import Presentation
except ImportError:
    print("Error: python-pptx is required. Install with: pip install python-pptx", file=sys.stderr)
    sys.exit(1)


def has_formatting(run):
    """Check if a run has explicit formatting (font properties set)."""
    rPr = run._r.find(
        '{http://schemas.openxmlformats.org/drawingml/2006/main}rPr'
    )
    return rPr is not None and len(rPr.attrib) > 0


def extract_paragraphs(text_frame):
    """Extract paragraph and run data from a text frame."""
    paragraphs = []
    for paragraph in text_frame.paragraphs:
        full_text = paragraph.text
        if not full_text.strip():
            continue
        runs = []
        for run_idx, run in enumerate(paragraph.runs):
            if not run.text:
                continue
            runs.append({
                "index": run_idx,
                "text": run.text,
                "has_formatting": has_formatting(run),
            })
        if runs:
            paragraphs.append({
                "full_text": full_text,
                "runs": runs,
            })
    return paragraphs


def extract_table(table):
    """Extract text from a table shape."""
    rows = []
    for row_idx, row in enumerate(table.rows):
        cells = []
        for cell_idx, cell in enumerate(row.cells):
            paragraphs = extract_paragraphs(cell.text_frame)
            if paragraphs:
                cells.append({
                    "row": row_idx,
                    "col": cell_idx,
                    "paragraphs": paragraphs,
                })
        if cells:
            rows.append({"cells": cells})
    return rows


def extract_shapes(shapes, depth=0):
    """Recursively extract text from shapes including grouped shapes."""
    elements = []
    for shape_idx, shape in enumerate(shapes):
        # Group shape: recurse into children
        if shape.shape_type is not None and shape.shape_type == 6:  # MSO_SHAPE_TYPE.GROUP
            try:
                group_elements = extract_shapes(shape.shapes, depth + 1)
                if group_elements:
                    elements.append({
                        "type": "group",
                        "shape_name": shape.name,
                        "shape_index": shape_idx,
                        "elements": group_elements,
                    })
            except AttributeError:
                pass
            continue

        # Table
        if shape.has_table:
            rows = extract_table(shape.table)
            if rows:
                elements.append({
                    "type": "table",
                    "shape_name": shape.name,
                    "shape_index": shape_idx,
                    "rows": rows,
                })
            continue

        # Text frame
        if shape.has_text_frame:
            paragraphs = extract_paragraphs(shape.text_frame)
            if paragraphs:
                elements.append({
                    "type": "text_frame",
                    "shape_name": shape.name,
                    "shape_index": shape_idx,
                    "paragraphs": paragraphs,
                })

    return elements


def extract_notes(slide):
    """Extract speaker notes from a slide."""
    if not slide.has_notes_slide:
        return []
    notes_frame = slide.notes_slide.notes_text_frame
    return extract_paragraphs(notes_frame)


def extract_pptx(pptx_path, target_lang="ko"):
    """Extract all text from a PPTX file at run level."""
    prs = Presentation(pptx_path)
    slides_data = []

    for slide_idx, slide in enumerate(prs.slides):
        elements = extract_shapes(slide.shapes)
        notes = extract_notes(slide)

        slide_data = {
            "slide_number": slide_idx + 1,
            "elements": elements,
        }
        if notes:
            slide_data["notes"] = notes

        slides_data.append(slide_data)

    return {
        "source_file": os.path.basename(pptx_path),
        "target_language": target_lang,
        "total_slides": len(slides_data),
        "slides": slides_data,
    }


def main():
    parser = argparse.ArgumentParser(
        description="Extract text from PPTX at run-level for translation"
    )
    parser.add_argument("input", help="Path to input PPTX file")
    parser.add_argument(
        "--per-slide",
        action="store_true",
        help="Output one JSON file per slide (for parallel translation)",
    )
    parser.add_argument(
        "--output-dir",
        default=None,
        help="Output directory (default: same as input file)",
    )
    parser.add_argument(
        "--target-lang",
        default="ko",
        help="Target language code (default: ko)",
    )
    args = parser.parse_args()

    pptx_path = os.path.abspath(args.input)
    if not os.path.exists(pptx_path):
        print(f"Error: File not found: {pptx_path}", file=sys.stderr)
        sys.exit(1)

    output_dir = args.output_dir or os.path.dirname(pptx_path)
    os.makedirs(output_dir, exist_ok=True)

    stem = Path(pptx_path).stem

    print(f"Extracting text from: {pptx_path}")
    data = extract_pptx(pptx_path, args.target_lang)

    if args.per_slide:
        # Output one JSON per slide
        for slide in data["slides"]:
            slide_num = slide["slide_number"]
            slide_file = {
                "source_file": data["source_file"],
                "target_language": data["target_language"],
                "total_slides": data["total_slides"],
                "slide": slide,
            }
            out_path = os.path.join(output_dir, f"{stem}_slide_{slide_num}.json")
            with open(out_path, "w", encoding="utf-8") as f:
                json.dump(slide_file, f, ensure_ascii=False, indent=2)
            print(f"  Slide {slide_num} -> {out_path}")
        print(f"Done: {data['total_slides']} slide files written to {output_dir}")
    else:
        # Single JSON output
        out_path = os.path.join(output_dir, f"{stem}_text.json")
        with open(out_path, "w", encoding="utf-8") as f:
            json.dump(data, f, ensure_ascii=False, indent=2)
        print(f"Done: {out_path}")
        print(f"  Total slides: {data['total_slides']}")

        # Summary stats
        total_elements = sum(len(s["elements"]) for s in data["slides"])
        total_runs = 0
        for s in data["slides"]:
            for el in s["elements"]:
                if el["type"] == "text_frame":
                    for p in el.get("paragraphs", []):
                        total_runs += len(p.get("runs", []))
                elif el["type"] == "table":
                    for row in el.get("rows", []):
                        for cell in row.get("cells", []):
                            for p in cell.get("paragraphs", []):
                                total_runs += len(p.get("runs", []))
                elif el["type"] == "group":
                    # Count runs in group recursively
                    def count_group_runs(elements):
                        count = 0
                        for e in elements:
                            if e["type"] == "text_frame":
                                for p in e.get("paragraphs", []):
                                    count += len(p.get("runs", []))
                            elif e["type"] == "table":
                                for row in e.get("rows", []):
                                    for cell in row.get("cells", []):
                                        for p in cell.get("paragraphs", []):
                                            count += len(p.get("runs", []))
                            elif e["type"] == "group":
                                count += count_group_runs(e.get("elements", []))
                        return count
                    total_runs += count_group_runs(el.get("elements", []))

        print(f"  Total text elements: {total_elements}")
        print(f"  Total runs: {total_runs}")


if __name__ == "__main__":
    main()

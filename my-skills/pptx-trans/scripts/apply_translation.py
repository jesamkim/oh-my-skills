#!/usr/bin/env python3
"""Apply translated JSON back to PPTX preserving all formatting.

Usage:
    python apply_translation.py input.pptx translated.json [--prefix KO_] [--output-dir DIR]
    python apply_translation.py input.pptx --per-slide [--prefix KO_] [--output-dir DIR]
    python apply_translation.py input.pptx translated.json --ea-font "Nanum Gothic"

Key principle: Only run.text is modified. Font objects are never touched.
When --ea-font is specified, East Asian font is set for runs containing CJK characters,
ensuring proper rendering when the original font (e.g., Amazon Ember) lacks CJK glyphs.
"""

import argparse
import glob
import json
import os
import re
import sys
from pathlib import Path

try:
    from pptx import Presentation
    from pptx.oxml.ns import qn
    from lxml import etree
except ImportError:
    print("Error: python-pptx is required. Install with: pip install python-pptx", file=sys.stderr)
    sys.exit(1)

# CJK Unicode ranges for detecting Korean/Japanese/Chinese characters
_CJK_RANGES = [
    (0xAC00, 0xD7AF),   # Hangul Syllables
    (0x1100, 0x11FF),   # Hangul Jamo
    (0x3130, 0x318F),   # Hangul Compatibility Jamo
    (0xA960, 0xA97F),   # Hangul Jamo Extended-A
    (0xD7B0, 0xD7FF),   # Hangul Jamo Extended-B
    (0x4E00, 0x9FFF),   # CJK Unified Ideographs
    (0x3040, 0x309F),   # Hiragana
    (0x30A0, 0x30FF),   # Katakana
]


def contains_cjk(text):
    """Check if text contains CJK (Korean/Japanese/Chinese) characters."""
    for ch in text:
        cp = ord(ch)
        for start, end in _CJK_RANGES:
            if start <= cp <= end:
                return True
    return False


def set_ea_font(run, font_name):
    """Set East Asian font for a run via XML, preserving all other properties."""
    rPr = run._r.get_or_add_rPr()
    ea = rPr.find(qn('a:ea'))
    if ea is None:
        ea = etree.SubElement(rPr, qn('a:ea'))
    ea.set('typeface', font_name)


# Korean font fallback order
_KOREAN_FONT_CANDIDATES = [
    "Nanum Gothic",
    "NanumGothic",
    "Apple SD Gothic Neo",    # macOS built-in
    "Malgun Gothic",          # Windows built-in (맑은 고딕)
    "맑은 고딕",
    "나눔고딕",
]


def _check_font_available(font_name):
    """Check if a font is available on the system."""
    import subprocess
    system = os.uname().sysname
    try:
        if system == "Darwin":  # macOS
            result = subprocess.run(
                ["system_profiler", "SPFontsDataType"],
                capture_output=True, text=True, timeout=10,
            )
            return font_name.lower() in result.stdout.lower()
        elif system == "Linux":
            result = subprocess.run(
                ["fc-list", f":family={font_name}"],
                capture_output=True, text=True, timeout=10,
            )
            return bool(result.stdout.strip())
        else:  # Windows or unknown
            return True  # Assume available, PowerPoint will handle fallback
    except (subprocess.TimeoutExpired, FileNotFoundError):
        return True  # If we can't check, assume available


def resolve_ea_font(requested_font):
    """Resolve East Asian font: check if requested font exists, try fallbacks."""
    if _check_font_available(requested_font):
        return requested_font

    print(f"  Font '{requested_font}' not found, trying fallbacks...")
    for candidate in _KOREAN_FONT_CANDIDATES:
        if candidate != requested_font and _check_font_available(candidate):
            print(f"  Using fallback font: {candidate}")
            return candidate

    print("  Warning: No Korean font found on system, skipping EA font setting")
    return None


def enable_autofit(text_frame):
    """Enable 'shrink text on overflow' for a text frame.

    Sets <a:normAutofit> in <a:bodyPr>, which tells PowerPoint to
    automatically reduce font size to fit text within the text box.
    This prevents text overflow when translated text is longer than
    the original.
    """
    bodyPr = text_frame._txBody.find(qn('a:bodyPr'))
    if bodyPr is None:
        return

    # Remove existing autofit settings
    for child_tag in ['a:noAutofit', 'a:spAutoFit']:
        existing = bodyPr.find(qn(child_tag))
        if existing is not None:
            bodyPr.remove(existing)

    # Add normAutofit if not already present
    norm = bodyPr.find(qn('a:normAutofit'))
    if norm is None:
        etree.SubElement(bodyPr, qn('a:normAutofit'))


def maybe_enable_autofit(text_frame, original_lengths, auto_fit):
    """Enable autofit if translated text is significantly longer than original.

    Compares total character count before and after translation.
    If the text grew by more than 20%, enable auto-shrink.
    """
    if not auto_fit:
        return

    new_total = sum(len(p.text) for p in text_frame.paragraphs)
    orig_total = sum(original_lengths)

    if orig_total > 0 and new_total > orig_total * 1.2:
        enable_autofit(text_frame)


def apply_paragraphs(text_frame, translated_paragraphs):
    """Apply translated text to paragraphs in a text frame, run by run."""
    applied = 0
    for t_para in translated_paragraphs:
        full_text = t_para.get("full_text", "")
        t_runs = t_para.get("runs", [])

        # Build a mapping of run index -> translated text
        run_map = {}
        for t_run in t_runs:
            idx = t_run["index"]
            # Use translated_text if present, otherwise keep original text
            translated = t_run.get("translated_text", t_run.get("text", ""))
            run_map[idx] = translated

        # Match paragraphs by full_text content
        for paragraph in text_frame.paragraphs:
            if paragraph.text.strip() == full_text.strip() or not full_text.strip():
                # Skip empty paragraphs
                if not paragraph.text.strip():
                    continue
                # Try exact match first
                if paragraph.text == full_text or paragraph.text.strip() == full_text.strip():
                    for run_idx, run in enumerate(paragraph.runs):
                        if run_idx in run_map:
                            run.text = run_map[run_idx]
                            applied += 1
                    break
    return applied


def apply_paragraphs_by_index(paragraphs_obj, translated_paragraphs, ea_font=None,
                              auto_fit=False, text_frame=None):
    """Apply translations using paragraph index matching within a text frame.

    Uses two-pointer approach to handle empty paragraphs in the PPTX that
    were skipped during extraction. Iterates PPTX paragraphs, skipping empty
    ones, and matches them sequentially with translated paragraphs.
    """
    applied = 0
    original_lengths = [len(p.text) for p in paragraphs_obj]

    t_idx = 0  # pointer into translated_paragraphs
    for paragraph in paragraphs_obj:
        if t_idx >= len(translated_paragraphs):
            break

        # Skip empty PPTX paragraphs without advancing translated pointer
        if not paragraph.text.strip():
            continue

        t_para = translated_paragraphs[t_idx]
        t_runs = t_para.get("runs", [])
        t_idx += 1

        run_map = {}
        for t_run in t_runs:
            idx = t_run["index"]
            translated = t_run.get("translated_text", t_run.get("text", ""))
            run_map[idx] = translated

        for run_idx, run in enumerate(paragraph.runs):
            if run_idx in run_map:
                run.text = run_map[run_idx]
                applied += 1
                # Set East Asian font for CJK characters
                if ea_font and contains_cjk(run.text):
                    set_ea_font(run, ea_font)

    # Enable autofit if text grew significantly
    if applied > 0 and text_frame is not None:
        maybe_enable_autofit(text_frame, original_lengths, auto_fit)

    return applied


def apply_to_shapes(shapes, elements, ea_font=None, auto_fit=False):
    """Apply translations to shapes, handling nested groups."""
    total_applied = 0

    for element in elements:
        el_type = element.get("type")
        shape_idx = element.get("shape_index")

        if shape_idx is not None and shape_idx < len(list(shapes)):
            shape = list(shapes)[shape_idx]

            if el_type == "text_frame" and shape.has_text_frame:
                paras = element.get("paragraphs", [])
                total_applied += apply_paragraphs_by_index(
                    shape.text_frame.paragraphs, paras, ea_font=ea_font,
                    auto_fit=auto_fit, text_frame=shape.text_frame,
                )

            elif el_type == "table" and shape.has_table:
                for row_data in element.get("rows", []):
                    for cell_data in row_data.get("cells", []):
                        row_i = cell_data.get("row", 0)
                        col_i = cell_data.get("col", 0)
                        try:
                            cell = shape.table.cell(row_i, col_i)
                            total_applied += apply_paragraphs_by_index(
                                cell.text_frame.paragraphs,
                                cell_data.get("paragraphs", []),
                                ea_font=ea_font,
                                auto_fit=auto_fit,
                                text_frame=cell.text_frame,
                            )
                        except (IndexError, KeyError):
                            pass

            elif el_type == "group":
                try:
                    total_applied += apply_to_shapes(
                        shape.shapes, element.get("elements", []),
                        ea_font=ea_font, auto_fit=auto_fit,
                    )
                except AttributeError:
                    pass

    return total_applied


def apply_notes(slide, notes_data, ea_font=None, auto_fit=False):
    """Apply translated text to speaker notes."""
    if not notes_data or not slide.has_notes_slide:
        return 0
    notes_frame = slide.notes_slide.notes_text_frame
    return apply_paragraphs_by_index(
        notes_frame.paragraphs, notes_data, ea_font=ea_font,
        auto_fit=auto_fit, text_frame=notes_frame,
    )


def repair_json(raw):
    """Attempt to repair JSON with unescaped quotes inside string values."""
    lines = raw.split('\n')
    fixed_lines = []
    for line in lines:
        match = re.match(
            r'^(\s*"(?:full_text|text|translated_text)":\s*)"(.+)"(,?)$', line
        )
        if match:
            prefix = match.group(1)
            value = match.group(2)
            suffix = match.group(3)
            # Temporarily protect already-escaped quotes
            value = value.replace('\\"', '\x00ESC\x00')
            # Escape remaining unescaped quotes
            value = value.replace('"', '\\"')
            # Restore previously escaped quotes
            value = value.replace('\x00ESC\x00', '\\"')
            fixed_lines.append(f'{prefix}"{value}"{suffix}')
        else:
            fixed_lines.append(line)
    return '\n'.join(fixed_lines)


def load_json_safe(fpath):
    """Load JSON with automatic repair for common agent writing errors."""
    with open(fpath, "r", encoding="utf-8") as f:
        raw = f.read()
    try:
        return json.loads(raw)
    except json.JSONDecodeError:
        repaired = repair_json(raw)
        try:
            data = json.loads(repaired)
            # Write back the repaired version
            with open(fpath, "w", encoding="utf-8") as f:
                json.dump(data, f, ensure_ascii=False, indent=2)
            print(f"  (auto-repaired JSON: {os.path.basename(fpath)})")
            return data
        except json.JSONDecodeError as e:
            print(f"Error: Cannot parse {fpath}: {e}", file=sys.stderr)
            sys.exit(1)


def load_translations(json_path=None, per_slide_dir=None, stem=None):
    """Load translation data from single file or per-slide files."""
    if json_path:
        data = load_json_safe(json_path)
        return data.get("slides", [data.get("slide")])

    if per_slide_dir and stem:
        slides = []
        pattern = os.path.join(per_slide_dir, f"{stem}_slide_*.json")
        files = sorted(glob.glob(pattern), key=lambda x: int(
            Path(x).stem.split("_slide_")[-1]
        ))
        for fpath in files:
            data = load_json_safe(fpath)
            slide_data = data.get("slide", data)
            slides.append(slide_data)
        return slides

    return []


def main():
    parser = argparse.ArgumentParser(
        description="Apply translated JSON back to PPTX preserving formatting"
    )
    parser.add_argument("input", help="Path to original PPTX file")
    parser.add_argument(
        "json_file",
        nargs="?",
        default=None,
        help="Path to translated JSON file (omit for --per-slide mode)",
    )
    parser.add_argument(
        "--per-slide",
        action="store_true",
        help="Read per-slide JSON files from output directory",
    )
    parser.add_argument(
        "--prefix",
        default="KO_",
        help="Output filename prefix (default: KO_)",
    )
    parser.add_argument(
        "--output-dir",
        default=None,
        help="Output directory (default: same as input file)",
    )
    parser.add_argument(
        "--ea-font",
        default="Nanum Gothic",
        help="East Asian font for CJK characters (default: Nanum Gothic). "
             "Use --no-ea-font to disable.",
    )
    parser.add_argument(
        "--no-ea-font",
        action="store_true",
        help="Disable East Asian font setting",
    )
    parser.add_argument(
        "--auto-fit",
        action="store_true",
        default=True,
        help="Enable auto-shrink for text frames where translated text is longer (default: on)",
    )
    parser.add_argument(
        "--no-auto-fit",
        action="store_true",
        help="Disable auto-shrink for overflowing text frames",
    )
    parser.add_argument(
        "--cleanup",
        action="store_true",
        default=True,
        help="Remove temporary per-slide JSON files after applying (default: on)",
    )
    parser.add_argument(
        "--no-cleanup",
        action="store_true",
        help="Keep temporary JSON files after applying",
    )
    args = parser.parse_args()

    pptx_path = os.path.abspath(args.input)
    if not os.path.exists(pptx_path):
        print(f"Error: File not found: {pptx_path}", file=sys.stderr)
        sys.exit(1)

    output_dir = args.output_dir or os.path.dirname(pptx_path)
    stem = Path(pptx_path).stem
    out_filename = f"{args.prefix}{Path(pptx_path).name}"
    out_path = os.path.join(output_dir, out_filename)

    # Load translations
    if args.per_slide:
        json_dir = args.output_dir or os.path.dirname(pptx_path)
        slides_data = load_translations(per_slide_dir=json_dir, stem=stem)
        if not slides_data:
            print(f"Error: No per-slide JSON files found for {stem}", file=sys.stderr)
            sys.exit(1)
        print(f"Loaded {len(slides_data)} per-slide translation files")
    elif args.json_file:
        json_path = os.path.abspath(args.json_file)
        if not os.path.exists(json_path):
            print(f"Error: JSON file not found: {json_path}", file=sys.stderr)
            sys.exit(1)
        slides_data = load_translations(json_path=json_path)
        print(f"Loaded translation from: {json_path}")
    else:
        print("Error: Provide a JSON file or use --per-slide mode", file=sys.stderr)
        sys.exit(1)

    # Determine East Asian font with system availability check
    ea_font = None if args.no_ea_font else args.ea_font
    if ea_font:
        ea_font = resolve_ea_font(ea_font)
        if ea_font:
            print(f"East Asian font: {ea_font}")
        else:
            print("East Asian font: disabled (no suitable font found)")

    # Determine auto-fit setting
    auto_fit = not args.no_auto_fit
    if auto_fit:
        print("Auto-fit: enabled (text frames will auto-shrink if text overflows)")

    # Open original PPTX
    prs = Presentation(pptx_path)
    total_applied = 0

    for slide_data in slides_data:
        if slide_data is None:
            continue
        slide_num = slide_data.get("slide_number", 0)
        if slide_num < 1 or slide_num > len(prs.slides):
            print(f"  Warning: Slide {slide_num} out of range, skipping")
            continue

        slide = prs.slides[slide_num - 1]
        elements = slide_data.get("elements", [])
        notes = slide_data.get("notes", [])

        applied = apply_to_shapes(slide.shapes, elements, ea_font=ea_font, auto_fit=auto_fit)
        applied += apply_notes(slide, notes, ea_font=ea_font, auto_fit=auto_fit)
        total_applied += applied

        if applied > 0:
            print(f"  Slide {slide_num}: {applied} runs translated")

    # Save
    os.makedirs(output_dir, exist_ok=True)
    prs.save(out_path)
    print(f"\nDone: {out_path}")
    print(f"  Total runs translated: {total_applied}")

    # Cleanup temporary JSON files
    do_cleanup = not args.no_cleanup
    if do_cleanup and args.per_slide:
        json_dir = args.output_dir or os.path.dirname(pptx_path)
        pattern = os.path.join(json_dir, f"{stem}_slide_*.json")
        json_files = glob.glob(pattern)
        if json_files:
            for f in json_files:
                os.remove(f)
            print(f"  Cleaned up {len(json_files)} temporary JSON files")
    elif do_cleanup and args.json_file:
        json_path = os.path.abspath(args.json_file)
        if os.path.exists(json_path):
            os.remove(json_path)
            print(f"  Cleaned up: {os.path.basename(json_path)}")


if __name__ == "__main__":
    main()
